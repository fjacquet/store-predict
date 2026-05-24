"""Native PowerPoint chart builders + Sankey picture for the PPTX report.

Charts that have a native PowerPoint equivalent (pie, column, bar) are added as
editable charts so the engineer can restyle them in PowerPoint. The Sankey flow
diagram has no native chart type, so it is embedded as an image rendered by the
shared matplotlib renderer in ``pdf_charts``.

Categories are aggregated by name (``_by_category``) so a workload that appears
in several DRR bands renders as a single, de-duplicated series — no repeated
``Database`` / ``File`` / ``VDI`` labels.
"""

from __future__ import annotations

from io import BytesIO
from typing import TYPE_CHECKING, Any

import i18n as _i18n
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.util import Pt

from store_predict.services.pdf_charts import render_sankey_png

if TYPE_CHECKING:
    from pptx.slide import Slide
    from pptx.util import Length

    from store_predict.pipeline.calculation import CalculationSummary

__all__ = [
    "CATEGORICAL_PALETTE_HEX",
    "add_before_after_bar",
    "add_capacity_bar",
    "add_drr_bar",
    "add_sankey_picture",
    "add_workload_pie",
]

# Midnight Executive categorical palette (matches the web app's ECharts donut).
CATEGORICAL_PALETTE_HEX = (
    "3245B7",  # navy primary
    "F9B935",  # gold
    "819AE9",  # light navy
    "4AA342",  # green
    "EF8700",  # orange
    "1E2761",  # deep navy
    "B0C2F9",  # ice
    "DF202E",  # red
    "2CC6B0",  # teal
    "64748B",  # slate
)
_PALETTE = tuple(RGBColor.from_string(h) for h in CATEGORICAL_PALETTE_HEX)
_NAVY = RGBColor.from_string("3245B7")
_DEEP = RGBColor.from_string("1E2761")
_LIGHT = RGBColor.from_string("819AE9")
_MUTED = RGBColor.from_string("64748B")


def t(key: str, **kwargs: object) -> str:
    """Translate via raw python-i18n (process-global locale set by the caller).

    Mirrors ``pptx_report.t``: these builders run inside ``run.io_bound`` where the
    tab-scoped wrapper would ignore the requested locale.
    """
    return str(_i18n.t(key, **kwargs))


def _by_category(summary: CalculationSummary) -> list[tuple[str, float, float, float]]:
    """Aggregate workload groups by category name, sorted by provisioned desc.

    Returns ``(name, provisioned_gib, required_gib, drr)`` tuples. Merges the
    per-DRR splits (e.g. two ``Database`` rows) into one entry per category so
    charts don't show duplicate labels.
    """
    totals: dict[str, list[float]] = {}
    order: list[str] = []
    for grp in summary.workload_groups:
        if grp.category not in totals:
            totals[grp.category] = [0.0, 0.0]
            order.append(grp.category)
        totals[grp.category][0] += grp.total_provisioned_mib
        totals[grp.category][1] += grp.total_required_mib
    rows = [
        (
            name,
            totals[name][0] / 1024,
            totals[name][1] / 1024,
            totals[name][0] / totals[name][1] if totals[name][1] else 0.0,
        )
        for name in order
    ]
    rows.sort(key=lambda r: r[1], reverse=True)
    return rows


def _style_axes(chart: Any, *, size: int = 9) -> None:
    """Shrink axis tick labels and mute gridlines for a cleaner look."""
    for axis_name in ("category_axis", "value_axis"):
        axis = getattr(chart, axis_name, None)
        if axis is None:
            continue
        try:
            axis.tick_labels.font.size = Pt(size)
            axis.tick_labels.font.color.rgb = _MUTED
        except (ValueError, NotImplementedError):
            pass


def _enable_data_labels(
    chart: Any, *, num_format: str, size: int = 9, position: XL_LABEL_POSITION | None = None
) -> None:
    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.number_format = num_format
    dl.number_format_is_linked = False
    dl.font.size = Pt(size)
    if position is not None:
        dl.position = position


def add_workload_pie(slide: Slide, summary: CalculationSummary, x: Length, y: Length, cx: Length, cy: Length) -> None:
    """Add an editable donut of provisioned capacity per workload category."""
    rows = _by_category(summary)
    if not rows:
        return
    data = CategoryChartData()
    data.categories = [r[0] for r in rows]
    data.add_series("GiB", tuple(r[1] for r in rows))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, x, y, cx, cy, data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(9)
    points = chart.series[0].points
    for idx, point in enumerate(points):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = _PALETTE[idx % len(_PALETTE)]
    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.show_percentage = True
    dl.show_value = False
    dl.number_format = "0%"
    dl.number_format_is_linked = False
    dl.font.size = Pt(8)
    dl.font.color.rgb = RGBColor.from_string("FFFFFF")


def add_capacity_bar(slide: Slide, summary: CalculationSummary, x: Length, y: Length, cx: Length, cy: Length) -> None:
    """Add an editable horizontal bar of provisioned capacity (GiB) per category."""
    rows = _by_category(summary)
    if not rows:
        return
    # BAR_CLUSTERED plots the first category at the bottom; reverse so the
    # largest category appears at the top.
    rows = list(reversed(rows))
    data = CategoryChartData()
    data.categories = [r[0] for r in rows]
    data.add_series(t("pdf.table_provisioned"), tuple(r[1] for r in rows))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, data).chart
    chart.has_legend = False
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = _NAVY
    _style_axes(chart)
    _enable_data_labels(chart, num_format="#,##0", size=9, position=XL_LABEL_POSITION.OUTSIDE_END)


def add_drr_bar(slide: Slide, summary: CalculationSummary, x: Length, y: Length, cx: Length, cy: Length) -> None:
    """Add an editable column chart of DRR per workload category."""
    rows = _by_category(summary)
    if not rows:
        return
    data = CategoryChartData()
    data.categories = [r[0] for r in rows]
    data.add_series("DRR", tuple(round(r[3], 2) for r in rows))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, data).chart
    chart.has_legend = False
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = _NAVY
    _style_axes(chart)
    _enable_data_labels(chart, num_format="0.0", size=9, position=XL_LABEL_POSITION.OUTSIDE_END)


def add_before_after_bar(
    slide: Slide, summary: CalculationSummary, x: Length, y: Length, cx: Length, cy: Length
) -> None:
    """Add an editable two-series column chart: provisioned vs required GiB per category."""
    rows = _by_category(summary)
    if not rows:
        return
    data = CategoryChartData()
    data.categories = [r[0] for r in rows]
    data.add_series(t("pptx.provisioned"), tuple(r[1] for r in rows))
    data.add_series(t("pptx.required"), tuple(r[2] for r in rows))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(10)
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = _NAVY
    chart.series[1].format.fill.solid()
    chart.series[1].format.fill.fore_color.rgb = _LIGHT
    _style_axes(chart)


def add_sankey_picture(slide: Slide, summary: CalculationSummary, x: Length, y: Length, cx: Length, cy: Length) -> None:
    """Embed the provisioned→required Sankey as a picture. No-op when there is no data."""
    # width_pt/height_pt only set the render aspect ratio; the picture is sized to (cx, cy).
    png = render_sankey_png(summary, width_pt=640, height_pt=240)
    if png is None:
        return
    slide.shapes.add_picture(BytesIO(png), x, y, width=cx, height=cy)
