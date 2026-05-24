"""PowerPoint report generator for StorePredict sizing decks.

Produces a branded, editable .pptx from a CalculationSummary: a concise
customer-facing pitch deck followed by a technical appendix. The visual language
matches the web app and the sibling vAtlas deck — the "Midnight Executive"
palette (navy + gold on a light canvas), white left-accent stat cards packed
into grids, a navy header band with a gold rule, and a footer on every slide.
Charts with a native PowerPoint equivalent (donut, column, bar) are editable;
the Sankey is an embedded image.
"""

from __future__ import annotations

from datetime import UTC, datetime
from io import BytesIO
from typing import TYPE_CHECKING, Any

import i18n as _i18n
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from store_predict.services import pptx_charts
from store_predict.services.pdf_report import _layout_metric_rows

if TYPE_CHECKING:
    from pptx.slide import Slide

    from store_predict.pipeline.calculation import CalculationSummary
    from store_predict.pipeline.health_checks import HealthCheckResult

__all__ = ["generate_report_pptx"]


def t(key: str, **kwargs: object) -> str:
    """Translate via raw python-i18n using the process-global locale.

    Deliberately NOT the tab-scoped ``store_predict.i18n.t`` wrapper: this module
    runs inside ``run.io_bound`` (a worker thread with no NiceGUI request context),
    where that wrapper's ``get_locale()`` falls back to the default and would ignore
    the ``locale`` argument. ``generate_report_pptx`` sets the locale once via
    ``_i18n.set("locale", locale)``; this mirrors ``excel_report``.
    """
    return str(_i18n.t(key, **kwargs))


# --- Geometry (16:9) ---------------------------------------------------------
_SLIDE_W = Inches(13.333)
_SLIDE_H = Inches(7.5)
_MARGIN = Inches(0.55)
_CONTENT_W = Inches(13.333 - 2 * 0.55)
_BAND_H = Inches(0.95)
_BODY_TOP = Inches(1.5)
_FOOTER_TOP = Inches(7.08)
_GRID_GAP = Inches(0.3)
_GRID_ROW_GAP = Inches(0.35)
_BLANK_LAYOUT = 6  # "Blank" layout in the default template

# --- Midnight Executive palette (matches the web app + vAtlas) ---------------
_PRIMARY = RGBColor.from_string("3245B7")  # navy primary — bands, hero bg, accent
_DEEP = RGBColor.from_string("1E2761")  # deep navy — title-slide tiles
_GOLD = RGBColor.from_string("F9B935")  # gold accent — rules, highlight
_CANVAS = RGBColor.from_string("F8FAFC")  # light content background
_INK = RGBColor.from_string("0F172A")  # near-black text / big numbers
_MUTED = RGBColor.from_string("64748B")  # small labels
_LINE = RGBColor.from_string("E2E8F0")  # card borders
_WHITE = RGBColor.from_string("FFFFFF")
_ICE = RGBColor.from_string("B0C2F9")  # light-blue on navy
_RED = RGBColor.from_string("DF202E")
_ORANGE = RGBColor.from_string("EF8700")
_FONT = "Arial"  # standard deck font (text, labels, headings, tables)
_MONO = "Consolas"  # monospace for numerals (KPI/stat values) — the on-screen idiom

_DATE = datetime.now(tz=UTC).strftime("%Y-%m-%d")


# --- Formatting --------------------------------------------------------------
def _compact(mib: float) -> str:
    """Single-unit storage label (no wrapping): GiB under 1 TiB, else TiB."""
    gib = mib / 1024.0
    if gib >= 1024.0:
        return f"{gib / 1024.0:.1f} TiB"
    return f"{gib:.0f} GiB"


# --- Primitive builders ------------------------------------------------------
def _new_slide(prs: Any, *, dark: bool = False) -> Slide:
    slide = prs.slides.add_slide(prs.slide_layouts[_BLANK_LAYOUT])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _PRIMARY if dark else _CANVAS
    return slide


def _no_line_no_shadow(shape: Any) -> None:
    shape.line.fill.background()
    shape.shadow.inherit = False


def _apply_font(run: Any, *, size: int, bold: bool = False, color: RGBColor = _INK, mono: bool = False) -> None:
    run.font.name = _MONO if mono else _FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_header_band(slide: Slide, heading: str) -> None:
    """Navy band across the top with a white heading and a gold bottom rule."""
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), _SLIDE_W, _BAND_H)
    band.fill.solid()
    band.fill.fore_color.rgb = _PRIMARY
    _no_line_no_shadow(band)
    tf = band.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = _MARGIN
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = heading
    _apply_font(run, size=26, bold=True, color=_WHITE)
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), _BAND_H, _SLIDE_W, Inches(0.05))
    rule.fill.solid()
    rule.fill.fore_color.rgb = _GOLD
    _no_line_no_shadow(rule)


def _add_footer(slide: Slide, project_name: str, *, on_dark: bool = False) -> None:
    parts = ["StorePredict"]
    if project_name:
        parts.append(project_name)
    parts.append(_DATE)
    box = slide.shapes.add_textbox(_MARGIN, _FOOTER_TOP, _CONTENT_W, Inches(0.3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "  ·  ".join(parts)
    _apply_font(run, size=9, color=_ICE if on_dark else _MUTED)


def _add_text(
    slide: Slide,
    text: str,
    left: Any,
    top: Any,
    width: Any,
    height: Any,
    *,
    size: int = 18,
    bold: bool = False,
    color: RGBColor = _INK,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _apply_font(run, size=size, bold=bold, color=color)


def _stat_card(
    slide: Slide,
    label: str,
    value: str,
    left: Any,
    top: Any,
    width: Any,
    height: Any,
    *,
    accent: RGBColor = _PRIMARY,
    dark: bool = False,
    value_pt: int = 22,
) -> None:
    """A card with a small label, a large value, and a colored left-accent bar.

    ``dark`` renders a deep-navy card (title slide); otherwise a white card with a
    subtle border on the light canvas.
    """
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = _DEEP if dark else _WHITE
    if dark:
        _no_line_no_shadow(card)
    else:
        card.line.color.rgb = _LINE
        card.line.width = Pt(0.75)
        card.shadow.inherit = False
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.09), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    _no_line_no_shadow(bar)

    tf = card.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.24)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.12)
    p_label = tf.paragraphs[0]
    r_label = p_label.add_run()
    r_label.text = label
    _apply_font(r_label, size=11, color=_ICE if dark else _MUTED)
    p_value = tf.add_paragraph()
    r_value = p_value.add_run()
    r_value.text = value
    _apply_font(r_value, size=value_pt, bold=True, color=_GOLD if dark else _INK, mono=True)


def _card_grid(
    slide: Slide,
    items: list[tuple[str, str, RGBColor]],
    *,
    top: Any,
    n_cols: int,
    card_h: Any,
    row_gap: Any = _GRID_ROW_GAP,
    gap: Any = _GRID_GAP,
) -> None:
    """Lay out ``(label, value, accent)`` items as a card grid filling the width."""
    card_w = Inches((13.333 - 2 * 0.55 - (n_cols - 1) * 0.3) / n_cols)
    for i, (label, value, accent) in enumerate(items):
        col = i % n_cols
        row = i // n_cols
        left = Inches(0.55 + col * (card_w.inches + gap.inches))
        y = Inches(top.inches + row * (card_h.inches + row_gap.inches))
        _stat_card(slide, label, value, left, y, card_w, card_h, accent=accent)


# --- Slides ------------------------------------------------------------------
def _slide_title(prs: Any, summary: CalculationSummary, project_name: str, company_logo_bytes: bytes | None) -> None:
    slide = _new_slide(prs, dark=True)
    title = project_name or t("pdf.report_title")
    _add_text(slide, "STOREPREDICT", _MARGIN, Inches(1.35), _CONTENT_W, Inches(0.4), size=14, bold=True, color=_GOLD)
    _add_text(slide, title, _MARGIN, Inches(1.95), _CONTENT_W, Inches(1.3), size=40, bold=True, color=_WHITE)
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _MARGIN, Inches(3.35), Inches(2.6), Inches(0.06))
    rule.fill.solid()
    rule.fill.fore_color.rgb = _GOLD
    _no_line_no_shadow(rule)
    subtitle = f"{t('pdf.report_title')}  ·  {_DATE}"
    _add_text(slide, subtitle, _MARGIN, Inches(3.6), _CONTENT_W, Inches(0.5), size=16, color=_ICE)

    tiles = [
        (t("stats.total_vms"), f"{summary.total_vms:,}", _GOLD),
        (t("stats.total_provisioned"), _compact(summary.total_provisioned_mib), _GOLD),
        (t("stats.required_capacity"), _compact(summary.total_required_mib), _GOLD),
        (t("stats.avg_drr"), f"{summary.weighted_avg_drr:.1f}x", _GOLD),
    ]
    card_w = Inches((13.333 - 2 * 0.55 - 3 * 0.3) / 4)
    for i, (label, value, accent) in enumerate(tiles):
        left = Inches(0.55 + i * (card_w.inches + 0.3))
        _stat_card(slide, label, value, left, Inches(4.7), card_w, Inches(1.5), accent=accent, dark=True, value_pt=24)

    if company_logo_bytes:
        slide.shapes.add_picture(BytesIO(company_logo_bytes), Inches(11.0), Inches(0.45), height=Inches(0.7))
    _add_footer(slide, project_name, on_dark=True)


def _slide_exec_summary(prs: Any, summary: CalculationSummary) -> None:
    slide = _new_slide(prs)
    _add_header_band(slide, t("pptx.exec_summary_heading"))
    items: list[tuple[str, str, RGBColor]] = [
        (t("stats.total_vms"), f"{summary.total_vms:,}", _PRIMARY),
        (t("stats.total_provisioned"), _compact(summary.total_provisioned_mib), _PRIMARY),
        (t("stats.total_in_use"), _compact(summary.total_in_use_mib), _PRIMARY),
        (t("stats.required_capacity"), _compact(summary.total_required_mib), _GOLD),
        (t("stats.avg_drr"), f"{summary.weighted_avg_drr:.1f}x", _PRIMARY),
        (t("stats.avg_storage"), _compact(summary.avg_vm_size_mib), _PRIMARY),
    ]
    # CPU/RAM columns are absent in many RVTools/LiveOptics exports — only show
    # those cards when there's data, so the deck never displays bare zeros.
    if summary.total_cpus > 0:
        items.append((t("stats.total_cpus"), f"{summary.total_cpus:,}", _PRIMARY))
    if summary.total_memory_mib > 0:
        items.append((t("stats.total_memory"), _compact(summary.total_memory_mib), _PRIMARY))
    n_cols = 4 if len(items) >= 7 else 3
    _card_grid(slide, items, top=Inches(1.85), n_cols=n_cols, card_h=Inches(1.75), row_gap=Inches(0.5))
    _add_footer(slide, "")


def _slide_drr_story(prs: Any, summary: CalculationSummary) -> None:
    slide = _new_slide(prs)
    _add_header_band(slide, t("pptx.drr_story_heading"))
    # Left column: DRR hero + provisioned + required stacked cards.
    col_w = Inches(3.5)
    _stat_card(
        slide,
        t("stats.avg_drr"),
        f"{summary.weighted_avg_drr:.1f}x",
        _MARGIN,
        Inches(1.7),
        col_w,
        Inches(1.4),
        accent=_GOLD,
        value_pt=30,
    )
    _stat_card(
        slide,
        t("stats.total_provisioned"),
        _compact(summary.total_provisioned_mib),
        _MARGIN,
        Inches(3.35),
        col_w,
        Inches(1.3),
        accent=_PRIMARY,
    )
    _stat_card(
        slide,
        t("stats.required_capacity"),
        _compact(summary.total_required_mib),
        _MARGIN,
        Inches(4.85),
        col_w,
        Inches(1.3),
        accent=_PRIMARY,
    )
    # Right: before/after bar chart.
    pptx_charts.add_before_after_bar(slide, summary, Inches(4.4), Inches(1.55), Inches(8.35), Inches(5.2))
    _add_footer(slide, "")


def _slide_workload_mix(prs: Any, summary: CalculationSummary) -> None:
    slide = _new_slide(prs)
    _add_header_band(slide, t("pptx.workload_mix_heading"))
    pptx_charts.add_capacity_bar(slide, summary, _MARGIN, _BODY_TOP, _CONTENT_W, Inches(5.3))
    _add_footer(slide, "")


def _slide_recommendation(prs: Any, summary: CalculationSummary, health_result: HealthCheckResult | None) -> None:
    from store_predict.pipeline.health_checks import Severity

    slide = _new_slide(prs)
    _add_header_band(slide, t("pptx.recommendation_heading"))
    row1: list[tuple[str, str, RGBColor]] = [
        (t("stats.required_capacity"), _compact(summary.total_required_mib), _GOLD),
        (t("stats.total_vms"), f"{summary.total_vms:,}", _PRIMARY),
        (t("stats.avg_drr"), f"{summary.weighted_avg_drr:.1f}x", _PRIMARY),
    ]
    _card_grid(slide, row1, top=Inches(1.8), n_cols=3, card_h=Inches(1.8))
    if health_result is not None and health_result.has_data and health_result.findings:
        counts = {sev: sum(1 for f in health_result.findings if f.severity == sev) for sev in Severity}
        sev_cards: list[tuple[str, str, RGBColor]] = [
            (t("pdf.findings_severity_critical"), str(counts.get(Severity.CRITICAL, 0)), _RED),
            (t("pdf.findings_severity_warning"), str(counts.get(Severity.WARNING, 0)), _ORANGE),
            (t("pdf.findings_severity_info"), str(counts.get(Severity.INFO, 0)), _PRIMARY),
        ]
        _card_grid(slide, sev_cards, top=Inches(4.05), n_cols=3, card_h=Inches(1.6))
    _add_footer(slide, "")


def _add_table(
    slide: Slide,
    rows: list[list[str]],
    left: Any,
    top: Any,
    width: Any,
    height: Any,
    *,
    numeric_from_col: int = 1,
    col_widths: list[float] | None = None,
) -> None:
    """Add a styled table: navy header, zebra body rows, right-aligned numerics."""
    n_rows = len(rows)
    n_cols = len(rows[0])
    gf = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = gf.table
    table.first_row = False
    table.horz_banding = False
    if col_widths is not None:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    for r in range(n_rows):
        for c in range(n_cols):
            cell = table.cell(r, c)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = _PRIMARY
            else:
                cell.fill.fore_color.rgb = _WHITE if r % 2 else RGBColor.from_string("EEF1FB")
            para = cell.text_frame.paragraphs[0]
            if c >= numeric_from_col:
                para.alignment = PP_ALIGN.RIGHT
            run = para.add_run()
            run.text = rows[r][c]
            _apply_font(run, size=12 if r == 0 else 11, bold=r == 0, color=_WHITE if r == 0 else _INK)


def _slide_breakdown_table(prs: Any, summary: CalculationSummary) -> None:
    slide = _new_slide(prs)
    _add_header_band(slide, t("report.breakdown_heading"))
    rows: list[list[str]] = [
        [
            t("pdf.table_category"),
            t("pdf.table_vms"),
            t("pdf.table_provisioned"),
            t("pdf.table_avg_drr"),
            t("pdf.table_required"),
        ]
    ]
    for grp in summary.workload_groups:
        rows.append(
            [
                grp.category,
                f"{grp.vm_count:,}",
                f"{grp.total_provisioned_mib / 1024:,.1f}",
                f"{grp.avg_drr:.2f}",
                f"{grp.total_required_mib / 1024:,.1f}",
            ]
        )
    rows.append(
        [
            t("pdf.table_total"),
            f"{summary.total_vms:,}",
            f"{summary.total_provisioned_mib / 1024:,.1f}",
            f"{summary.weighted_avg_drr:.2f}",
            f"{summary.total_required_mib / 1024:,.1f}",
        ]
    )
    _add_table(slide, rows, _MARGIN, _BODY_TOP, _CONTENT_W, Inches(5.4))
    _add_footer(slide, "")


def _slide_layout_strategies(prs: Any, summary: CalculationSummary) -> None:
    from store_predict.pipeline.layout_engine import generate_all_proposals

    proposals = generate_all_proposals(summary)
    slide = _new_slide(prs)
    _add_header_band(slide, t("pdf.layout_heading"))
    rows: list[list[str]] = [
        [t("layout_page.metric"), t("strategy.consolidation"), t("strategy.performance"), t("strategy.uniform")]
    ]
    for metric_key, c_val, p_val, u_val in _layout_metric_rows(proposals):
        rows.append([t(f"metrics.{metric_key}"), c_val, p_val, u_val])
    _add_table(slide, rows, _MARGIN, _BODY_TOP, _CONTENT_W, Inches(5.45))
    _add_footer(slide, "")


def _slide_findings(prs: Any, health_result: HealthCheckResult) -> None:
    from store_predict.pipeline.health_checks import Severity

    slide = _new_slide(prs)
    _add_header_band(slide, t("pdf.findings_summary_heading"))
    order = {Severity.CRITICAL: 0, Severity.WARNING: 1, Severity.INFO: 2}
    findings = sorted(health_result.findings, key=lambda f: (order.get(f.severity, 3), f.check_id))
    rows: list[list[str]] = [
        [t("pdf.findings_col_severity"), t("pdf.findings_col_finding"), t("pdf.findings_col_count")]
    ]
    rows.extend([t(f"pdf.findings_severity_{f.severity}"), t(f.title), f"{f.affected_count:,}"] for f in findings)
    _add_table(
        slide,
        rows,
        _MARGIN,
        _BODY_TOP,
        _CONTENT_W,
        Inches(5.4),
        numeric_from_col=2,
        col_widths=[2.0, 8.233, 2.0],
    )
    _add_footer(slide, "")


def _slide_charts(prs: Any, summary: CalculationSummary) -> None:
    slide = _new_slide(prs)
    _add_header_band(slide, t("pdf.charts_heading"))
    pptx_charts.add_sankey_picture(slide, summary, _MARGIN, _BODY_TOP, _CONTENT_W, Inches(2.5))
    pptx_charts.add_workload_pie(slide, summary, _MARGIN, Inches(4.1), Inches(6.0), Inches(2.9))
    pptx_charts.add_drr_bar(slide, summary, Inches(6.85), Inches(4.1), Inches(6.0), Inches(2.9))
    _add_footer(slide, "")


def generate_report_pptx(
    summary: CalculationSummary,
    project_name: str,
    locale: str = "fr",
    company_logo_bytes: bytes | None = None,
    health_result: HealthCheckResult | None = None,
) -> bytes:
    """Generate a branded PPTX sizing deck and return raw bytes.

    Args:
        summary: Calculation results to render.
        project_name: Customer / project label for the title slide.
        locale: Language for deck labels (e.g. "fr" or "en"). Defaults to "fr".
        company_logo_bytes: Optional customer logo (already validated upstream).
        health_result: Optional health findings; drives the findings slide.

    Returns:
        The .pptx document as bytes.
    """
    # Set the process-global locale before any t() call. Safe within one call (this
    # function is synchronous — no coroutine interleaving), matching pdf_report/
    # excel_report. NOTE: not thread-safe across *concurrent* report generations in
    # different locales (run.io_bound uses a thread pool); benign in the single-
    # container deployment. A lock or per-call locale= would be needed for multi-user.
    _i18n.set("locale", locale)

    prs = Presentation()
    prs.slide_width = _SLIDE_W
    prs.slide_height = _SLIDE_H

    _slide_title(prs, summary, project_name, company_logo_bytes)
    _slide_exec_summary(prs, summary)
    _slide_drr_story(prs, summary)
    _slide_workload_mix(prs, summary)
    _slide_recommendation(prs, summary, health_result)

    # --- Appendix ---
    _slide_breakdown_table(prs, summary)
    if summary.total_vms > 0:
        _slide_layout_strategies(prs, summary)
    if health_result is not None and health_result.has_data and health_result.findings:
        _slide_findings(prs, health_result)
    _slide_charts(prs, summary)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
