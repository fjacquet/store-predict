"""Tests for the PPTX report generator service."""

from __future__ import annotations

from io import BytesIO

from pptx import Presentation

from store_predict.pipeline.calculation import (
    CalculationSummary,
    VMCalculation,
    WorkloadGroupResult,
)
from store_predict.pipeline.health_checks import HealthCheckResult, HealthFinding, Severity
from store_predict.services.pptx_report import generate_report_pptx


def _make_summary(
    categories: list[tuple[str, int, float, float]] | None = None,
) -> CalculationSummary:
    if categories is None:
        categories = []
    vm_calcs: list[VMCalculation] = []
    groups: list[WorkloadGroupResult] = []
    for cat_name, count, prov_mib, drr in categories:
        in_use = prov_mib * 0.6
        req = prov_mib / max(drr, 0.1)
        for i in range(count):
            vm_calcs.append(
                VMCalculation(
                    vm_name=f"{cat_name}-VM{i + 1}",
                    workload_category=cat_name,
                    provisioned_mib=prov_mib / count,
                    in_use_mib=in_use / count,
                    drr=drr,
                    required_mib=req / count,
                )
            )
        groups.append(
            WorkloadGroupResult(
                category=cat_name,
                vm_count=count,
                total_provisioned_mib=prov_mib,
                total_in_use_mib=in_use,
                avg_drr=drr,
                total_required_mib=req,
            )
        )
    total_prov = sum(g.total_provisioned_mib for g in groups)
    total_in_use = sum(g.total_in_use_mib for g in groups)
    total_req = sum(g.total_required_mib for g in groups)
    weighted = total_prov / total_req if total_req > 0 else 0.0
    return CalculationSummary(
        vm_calculations=vm_calcs,
        workload_groups=groups,
        total_vms=len(vm_calcs),
        total_provisioned_mib=total_prov,
        total_in_use_mib=total_in_use,
        total_required_mib=total_req,
        weighted_avg_drr=weighted,
        largest_vm_name=vm_calcs[0].vm_name if vm_calcs else "",
        largest_vm_provisioned_mib=vm_calcs[0].provisioned_mib if vm_calcs else 0.0,
    )


def _make_health_result() -> HealthCheckResult:
    findings = (
        HealthFinding(
            check_id="data_quality.missing_os",
            severity=Severity.WARNING,
            title="health.missing_os.title",
            detail="health.missing_os.detail",
            affected_count=3,
            affected_vms=("vm1", "vm2", "vm3"),
        ),
        HealthFinding(
            check_id="best_practice.tools_not_installed",
            severity=Severity.CRITICAL,
            title="health.tools_not_installed.title",
            detail="health.tools_not_installed.detail",
            affected_count=1,
            affected_vms=("vm1",),
        ),
    )
    return HealthCheckResult(findings=findings, total_vms_checked=10, has_data=True)


def _slide_text(prs: Presentation) -> str:
    """Concatenate all text from all slides for substring assertions."""
    chunks: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                chunks.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        chunks.append(cell.text)
    return "\n".join(chunks)


class TestPptxGeneratesBytes:
    def test_returns_pptx_magic_bytes(self) -> None:
        summary = _make_summary([("Database/Microsoft SQL", 3, 30720.0, 5.0)])
        result = generate_report_pptx(summary, "Test Project")
        assert isinstance(result, bytes)
        assert result[:4] == b"PK\x03\x04"  # .pptx is a zip container

    def test_opens_as_presentation_with_title(self) -> None:
        summary = _make_summary([("Database/Microsoft SQL", 3, 30720.0, 5.0)])
        prs = Presentation(BytesIO(generate_report_pptx(summary, "Acme Corp", locale="en")))
        assert len(prs.slides) >= 2
        assert "Acme Corp" in _slide_text(prs)


class TestMainDeck:
    def test_main_deck_has_charts_and_recommendation(self) -> None:
        summary = _make_summary([("Database/Microsoft SQL", 3, 30720.0, 5.0), ("Virtual Machines", 2, 10240.0, 5.0)])
        prs = Presentation(BytesIO(generate_report_pptx(summary, "X", locale="en")))
        # title + exec + drr-story + workload-mix + recommendation = 5 main slides
        assert len(prs.slides) >= 5
        # at least two native charts across the deck (before/after bar + workload pie)
        chart_count = sum(1 for slide in prs.slides for shape in slide.shapes if shape.has_chart)
        assert chart_count >= 2
        text = _slide_text(prs)
        assert "Recommendation" in text  # pptx.recommendation_heading (en)


class TestAppendix:
    def test_breakdown_table_present(self) -> None:
        summary = _make_summary([("Database/Microsoft SQL", 3, 30720.0, 5.0)])
        prs = Presentation(BytesIO(generate_report_pptx(summary, "X", locale="en")))
        assert any(shape.has_table for slide in prs.slides for shape in slide.shapes)
        assert "Category" in _slide_text(prs)  # pdf.table_category (en)

    def test_layout_slide_present_with_vms_absent_when_empty(self) -> None:
        with_vms = _make_summary([("Virtual Machines", 2, 10240.0, 5.0)])
        n_with = len(Presentation(BytesIO(generate_report_pptx(with_vms, "X", locale="en"))).slides)
        empty = _make_summary()  # total_vms == 0
        n_empty = len(Presentation(BytesIO(generate_report_pptx(empty, "X", locale="en"))).slides)
        assert n_with > n_empty

    def test_findings_slide_added_only_with_findings(self) -> None:
        summary = _make_summary([("Virtual Machines", 2, 10240.0, 5.0)])
        n_no = len(Presentation(BytesIO(generate_report_pptx(summary, "X", locale="en"))).slides)
        n_yes = len(
            Presentation(
                BytesIO(generate_report_pptx(summary, "X", locale="en", health_result=_make_health_result()))
            ).slides
        )
        assert n_yes == n_no + 1
