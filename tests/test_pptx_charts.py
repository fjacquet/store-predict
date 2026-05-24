"""Tests for the PPTX chart builders."""

from __future__ import annotations

from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from store_predict.pipeline.calculation import (
    CalculationSummary,
    VMCalculation,
    WorkloadGroupResult,
)
from store_predict.services import pptx_charts


def _make_summary() -> CalculationSummary:
    groups = [
        WorkloadGroupResult("Database/Microsoft SQL", 3, 30720.0, 18432.0, 5.0, 6144.0),
        WorkloadGroupResult("Virtual Machines", 2, 10240.0, 6144.0, 5.0, 2048.0),
    ]
    vm_calcs = [VMCalculation(f"VM-{i}", "Virtual Machines", 5120.0, 3072.0, 5.0, 1024.0) for i in range(5)]
    return CalculationSummary(
        vm_calculations=vm_calcs,
        workload_groups=groups,
        total_vms=5,
        total_provisioned_mib=40960.0,
        total_in_use_mib=24576.0,
        total_required_mib=8192.0,
        weighted_avg_drr=5.0,
    )


def _blank_slide():
    prs = Presentation()
    return prs.slides.add_slide(prs.slide_layouts[6])


def test_add_workload_pie_adds_doughnut_chart() -> None:
    slide = _blank_slide()
    pptx_charts.add_workload_pie(slide, _make_summary(), Inches(1), Inches(1), Inches(5), Inches(4))
    charts = [s.chart for s in slide.shapes if s.has_chart]
    assert len(charts) == 1
    assert charts[0].chart_type == XL_CHART_TYPE.DOUGHNUT


def test_add_capacity_bar_adds_bar_chart() -> None:
    slide = _blank_slide()
    pptx_charts.add_capacity_bar(slide, _make_summary(), Inches(1), Inches(1), Inches(5), Inches(4))
    charts = [s.chart for s in slide.shapes if s.has_chart]
    assert len(charts) == 1
    assert charts[0].chart_type == XL_CHART_TYPE.BAR_CLUSTERED


def test_add_drr_bar_adds_column_chart() -> None:
    slide = _blank_slide()
    pptx_charts.add_drr_bar(slide, _make_summary(), Inches(1), Inches(1), Inches(5), Inches(4))
    charts = [s.chart for s in slide.shapes if s.has_chart]
    assert len(charts) == 1
    assert charts[0].chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED


def test_add_before_after_bar_has_two_series() -> None:
    slide = _blank_slide()
    pptx_charts.add_before_after_bar(slide, _make_summary(), Inches(1), Inches(1), Inches(5), Inches(4))
    charts = [s.chart for s in slide.shapes if s.has_chart]
    assert len(charts) == 1
    assert len(charts[0].series) == 2


def test_chart_builders_noop_on_empty_summary() -> None:
    empty = CalculationSummary([], [], 0, 0.0, 0.0, 0.0, 0.0)
    slide = _blank_slide()
    pptx_charts.add_workload_pie(slide, empty, Inches(1), Inches(1), Inches(5), Inches(4))
    pptx_charts.add_capacity_bar(slide, empty, Inches(1), Inches(1), Inches(5), Inches(4))
    pptx_charts.add_drr_bar(slide, empty, Inches(1), Inches(1), Inches(5), Inches(4))
    pptx_charts.add_before_after_bar(slide, empty, Inches(1), Inches(1), Inches(5), Inches(4))
    assert not any(s.has_chart for s in slide.shapes)


def test_add_sankey_picture_adds_picture() -> None:
    slide = _blank_slide()
    pptx_charts.add_sankey_picture(slide, _make_summary(), Inches(1), Inches(1), Inches(8), Inches(3))
    pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pics) == 1


def test_add_sankey_picture_noop_on_empty() -> None:
    empty = CalculationSummary([], [], 0, 0.0, 0.0, 0.0, 0.0)
    slide = _blank_slide()
    pptx_charts.add_sankey_picture(slide, empty, Inches(1), Inches(1), Inches(8), Inches(3))
    assert not any(s.shape_type == MSO_SHAPE_TYPE.PICTURE for s in slide.shapes)
